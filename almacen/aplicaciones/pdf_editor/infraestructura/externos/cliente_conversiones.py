# -*- coding: utf-8 -*-
"""
Cliente de conversiones "premium" del Editor PDF.

Todas las funciones reciben bytes y devuelven bytes (o lista de tuplas), sin
tocar disco fuera de temporales propios: son reutilizables desde cualquier
interfaz (API, tareas programadas, scripts).

Conversiones soportadas:
- Office (Word/Excel/PowerPoint/ODF/RTF/HTML/TXT) → PDF: via Gotenberg
  (contenedor Docker local con LibreOffice, puerto 127.0.0.1:3000).
- PDF → Word (.docx): pdf2docx (mantiene layout, tablas e imágenes).
- PDF → Excel (.xlsx): pdfplumber detecta tablas y openpyxl arma el libro;
  si una página no tiene tablas se vuelca su texto línea a línea.
- PDF → PowerPoint (.pptx): cada página se renderiza a imagen (PyMuPDF)
  y se coloca en una diapositiva a tamaño completo (python-pptx).
- Numerar páginas, dividir por rangos y desbloquear (quitar contraseña).
"""

import io
import logging
import os
import tempfile

import fitz  # PyMuPDF

logger = logging.getLogger(__name__)

# URL del contenedor Gotenberg (Docker local). Si algún día se mueve a otra
# VM, basta cambiar esta constante o definir GOTENBERG_URL en el entorno.
GOTENBERG_URL = os.environ.get('GOTENBERG_URL', 'http://127.0.0.1:3000')

# Extensiones que Gotenberg/LibreOffice sabe convertir a PDF
EXTENSIONES_OFICINA = {
    '.doc', '.docx', '.odt', '.rtf', '.txt', '.html', '.htm',
    '.xls', '.xlsx', '.ods', '.csv',
    '.ppt', '.pptx', '.odp',
}


class ErrorConversion(Exception):
    """Error de conversión con mensaje apto para mostrar al usuario."""


# ============================================================
# EJECUCIÓN EN SUBPROCESO (conversiones CPU-intensivas)
# ============================================================

def en_subproceso(operacion, entradas, timeout=1500, params=None):
    """Ejecuta una operación pesada en un proceso aparte (conversor_cli.py).

    Por qué: pdf2docx/OpenCV/tesseract son CPU puro; dentro del worker eventlet
    congelan el event loop → gunicorn deja de recibir latidos y mata el worker
    a los 300 s → 504 en nginx y usuarios colgados. En subproceso, la espera es
    cooperativa (eventlet parchea subprocess): el worker sigue atendiendo.

    Args:
        operacion: pdf-a-word | pdf-a-excel | pdf-a-ppt | comparar | comprimir | ocr
        entradas: lista de bytes (1 archivo; comparar recibe 2)
        timeout: segundos máximos (por debajo de los 1800 s del nginx de /api/pdf/)
        params: dict de parámetros extra de la operación (calidad, idioma, página…)
    Returns:
        bytes del resultado (en 'ocr' es JSON codificado en UTF-8)
    """
    import json
    import subprocess
    import sys

    rutas = []
    fd, ruta_salida = tempfile.mkstemp(suffix='.out')
    os.close(fd)
    rutas.append(ruta_salida)
    try:
        rutas_entrada = []
        for contenido in entradas:
            fd, ruta = tempfile.mkstemp(suffix='.pdf')
            os.close(fd)
            rutas.append(ruta)
            rutas_entrada.append(ruta)
            with open(ruta, 'wb') as f:
                f.write(contenido)

        # gunicorn corre con PYTHONOPTIMIZE=2 (-OO); python-docx/pptx no lo toleran bien
        entorno = dict(os.environ)
        entorno.pop('PYTHONOPTIMIZE', None)

        cli = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'conversor_cli.py')
        try:
            r = subprocess.run([sys.executable, cli, operacion, ruta_salida,
                                json.dumps(params or {})] + rutas_entrada,
                               stdout=subprocess.DEVNULL, stderr=subprocess.PIPE,
                               timeout=timeout, env=entorno)
        except subprocess.TimeoutExpired:
            raise ErrorConversion('La conversión tardó más de %d minutos y se canceló. '
                                  'Intenta con un documento más pequeño.' % (timeout // 60))
        if r.returncode != 0:
            lineas = [l for l in r.stderr.decode('utf-8', 'replace').strip().split('\n') if l.strip()]
            # Antes se guardaban solo las 3 últimas líneas y el motivo de verdad se
            # perdía: pdf2docx descarta las páginas que le dan error una a una y al
            # final solo dice «No parsed pages». Ahora se guardan también las líneas
            # que explican POR QUÉ se descartó cada página (medido el 30-jul-2026,
            # con un fallo real que no se pudo diagnosticar por esto mismo).
            motivos = [l for l in lineas if 'Ignore page' in l or 'error' in l.lower()]
            logger.error('conversor_cli %s falló. Últimas líneas: %s | Motivos: %s',
                         operacion, lineas[-3:], motivos[:6])
            raise ErrorConversion(lineas[-1] if lineas else 'La conversión falló.')
        with open(ruta_salida, 'rb') as f:
            return f.read()
    finally:
        for ruta in rutas:
            try:
                os.remove(ruta)
            except OSError:
                pass


# ============================================================
# OFFICE → PDF (Gotenberg)
# ============================================================

def oficina_a_pdf(nombre_archivo, contenido, timeout=120):
    """Convierte un documento de oficina a PDF usando Gotenberg.

    Args:
        nombre_archivo: nombre original (la extensión decide el filtro de LibreOffice)
        contenido: bytes del archivo
    Returns:
        bytes del PDF
    """
    import requests

    ext = os.path.splitext(nombre_archivo or '')[1].lower()
    if ext not in EXTENSIONES_OFICINA:
        raise ErrorConversion(
            'Formato "%s" no admitido. Se aceptan: %s' % (ext or 'desconocido',
            ', '.join(sorted(EXTENSIONES_OFICINA))))

    try:
        resp = requests.post(
            GOTENBERG_URL + '/forms/libreoffice/convert',
            files={'files': (nombre_archivo, contenido)},
            timeout=timeout,
        )
    except requests.ConnectionError:
        raise ErrorConversion('El servicio de conversión (Gotenberg) no está disponible. '
                              'Verificar: docker ps | grep gotenberg')
    except requests.Timeout:
        raise ErrorConversion('La conversión tardó demasiado (más de %d s). '
                              'Intenta con un archivo más liviano.' % timeout)

    if resp.status_code != 200:
        logger.error('Gotenberg respondió %s: %s', resp.status_code, resp.text[:300])
        raise ErrorConversion('El convertidor no pudo procesar el archivo (código %d).' % resp.status_code)

    if not resp.content.startswith(b'%PDF'):
        raise ErrorConversion('El convertidor devolvió un resultado inválido.')
    return resp.content


# ============================================================
# PDF → WORD
# ============================================================

def pdf_a_word(contenido_pdf):
    """Convierte un PDF a .docx editable con pdf2docx (layout, tablas e imágenes).

    El resultado son párrafos y tablas de Word de verdad, editables: no imágenes de las
    páginas. Eso es lo que hace que la conversión cueste, porque hay que reconstruir la
    maquetación página a página.

    Se convierte **repartiendo las páginas entre varios procesos**: es lo que más pesa en
    la espera del usuario y el servidor tiene núcleos de sobra. Medido sobre un documento
    de 130 páginas en la VM 101 (24 núcleos): **36,1 s en serie → 11,1 s con 8 procesos**.
    Si el paralelo fallara por lo que sea, se reintenta en serie: más vale tardar que no
    entregar el archivo.
    """
    from pdf2docx import Converter

    fd_pdf, ruta_pdf = tempfile.mkstemp(suffix='.pdf')
    fd_docx, ruta_docx = tempfile.mkstemp(suffix='.docx')
    os.close(fd_pdf)
    os.close(fd_docx)
    try:
        with open(ruta_pdf, 'wb') as f:
            f.write(contenido_pdf)
        nucleos = max(1, min(8, (os.cpu_count() or 2) - 2))
        for kwargs in ({'multi_processing': True, 'cpu_count': nucleos}, {}):
            cv = Converter(ruta_pdf)
            try:
                cv.convert(ruta_docx, **kwargs)
                break
            except Exception as e:
                if not kwargs:          # el intento en serie era el último
                    raise
                logger.warning('pdf2docx en paralelo falló (%s): se reintenta en serie', e)
            finally:
                cv.close()
        with open(ruta_docx, 'rb') as f:
            return f.read()
    except Exception as e:
        logger.error('pdf2docx falló: %s', e)
        raise ErrorConversion('No se pudo convertir a Word: %s' % e)
    finally:
        for ruta in (ruta_pdf, ruta_docx):
            try:
                os.remove(ruta)
            except OSError:
                pass


# ============================================================
# PDF → EXCEL
# ============================================================

def _excel_lote(argumentos):
    """Lee un tramo de páginas con pdfplumber. Se ejecuta en un proceso aparte.

    Devuelve [(nº de página, 'tablas'|'texto', datos)] para que el proceso principal
    monte el libro. Tiene que ser una función de módulo: las anidadas no se pueden
    enviar a otro proceso.
    """
    contenido, desde, hasta = argumentos
    import pdfplumber
    salida = []
    with pdfplumber.open(io.BytesIO(contenido)) as pdf:
        for num in range(desde, min(hasta, len(pdf.pages))):
            pagina = pdf.pages[num]
            tablas = pagina.extract_tables()
            if tablas:
                salida.append((num + 1, 'tablas', tablas))
            else:
                texto = (pagina.extract_text() or '').strip()
                if texto:
                    salida.append((num + 1, 'texto', texto))
    return salida


def pdf_a_excel(contenido_pdf):
    """Extrae las tablas del PDF a un .xlsx (una hoja por página con contenido).

    Páginas con tablas → cada tabla se copia celda a celda.
    Páginas sin tablas → se vuelca el texto línea a línea (columna A).

    Las páginas se reparten entre varios procesos: buscar tablas es lo que cuesta y se
    hace página a página, así que escala bien. Medido con 130 páginas en la VM 101:
    **38,7 s en serie → unos 8 s con 8 procesos**. Si el reparto fallara, se hace en
    serie como siempre.
    """
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)  # se crean hojas solo para páginas con contenido
    hojas = 0

    try:
        import fitz as _fitz
        with _fitz.open(stream=contenido_pdf, filetype='pdf') as doc:
            total = len(doc)
        procesos = max(1, min(8, (os.cpu_count() or 2) - 2))
        paginas = []
        if total > 4 and procesos > 1:
            try:
                from concurrent.futures import ProcessPoolExecutor
                tam = max(1, -(-total // procesos))      # división hacia arriba
                lotes = [(contenido_pdf, i, i + tam) for i in range(0, total, tam)]
                with ProcessPoolExecutor(max_workers=procesos) as pool:
                    for parcial in pool.map(_excel_lote, lotes):
                        paginas.extend(parcial)
            except Exception as e:
                logger.warning('Excel en paralelo falló (%s): se hace en serie', e)
                paginas = []
        if not paginas:
            paginas = _excel_lote((contenido_pdf, 0, total))

        for num, clase, datos in sorted(paginas):
            if clase == 'tablas':
                ws = wb.create_sheet('Pág %d' % num)
                fila_actual = 1
                for tabla in datos:
                    for fila in tabla:
                        for col, valor in enumerate(fila, start=1):
                            ws.cell(row=fila_actual, column=col,
                                    value=(valor or '').strip() if isinstance(valor, str) else valor)
                        fila_actual += 1
                    fila_actual += 1  # renglón en blanco entre tablas
            else:
                ws = wb.create_sheet('Pág %d (texto)' % num)
                for i, linea in enumerate(datos.split('\n'), start=1):
                    ws.cell(row=i, column=1, value=linea)
            hojas += 1
    except Exception as e:
        logger.error('pdfplumber falló: %s', e)
        raise ErrorConversion('No se pudo leer el PDF para extraer tablas: %s' % e)

    if hojas == 0:
        raise ErrorConversion('El PDF no tiene tablas ni texto extraíble '
                              '(¿es un escaneado? Usa primero "Digitalizar y OCR").')

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ============================================================
# PDF → POWERPOINT
# ============================================================

def pdf_a_ppt(contenido_pdf, dpi=140):
    """Convierte el PDF a .pptx: una diapositiva por página (imagen a página completa)."""
    from pptx import Presentation
    from pptx.util import Emu

    EMU_POR_PUNTO = 12700  # 1 pt = 12700 EMU

    try:
        doc = fitz.open(stream=contenido_pdf, filetype='pdf')
    except Exception as e:
        raise ErrorConversion('PDF inválido: %s' % e)

    pres = Presentation()
    en_blanco = pres.slide_layouts[6]  # layout vacío

    try:
        for pagina in doc:
            ancho_emu = int(pagina.rect.width * EMU_POR_PUNTO)
            alto_emu = int(pagina.rect.height * EMU_POR_PUNTO)
            pres.slide_width = Emu(ancho_emu)
            pres.slide_height = Emu(alto_emu)

            pix = pagina.get_pixmap(dpi=dpi)
            slide = pres.slides.add_slide(en_blanco)
            slide.shapes.add_picture(io.BytesIO(pix.tobytes('png')), 0, 0,
                                     width=Emu(ancho_emu), height=Emu(alto_emu))
    finally:
        doc.close()

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


# ============================================================
# NUMERAR PÁGINAS
# ============================================================

def numerar_paginas(contenido_pdf, posicion='abajo-centro', formato='{n} de {total}',
                    tamano=10, desde=1):
    """Estampa el número de página en todas las páginas.

    Args:
        posicion: abajo-centro | abajo-derecha | abajo-izquierda | arriba-centro | arriba-derecha | arriba-izquierda
        formato: texto con {n} (número) y {total} (total de páginas)
        desde: número inicial (útil si el documento se anexa a otro)
    """
    try:
        doc = fitz.open(stream=contenido_pdf, filetype='pdf')
    except Exception as e:
        raise ErrorConversion('PDF inválido: %s' % e)

    total = len(doc)
    margen = 28  # pt desde el borde
    try:
        for i, pagina in enumerate(doc):
            texto = formato.replace('{n}', str(i + desde)).replace('{total}', str(total + desde - 1))
            ancho_texto = fitz.get_text_length(texto, fontsize=tamano)
            r = pagina.rect
            if 'izquierda' in posicion:
                x = margen
            elif 'derecha' in posicion:
                x = r.width - margen - ancho_texto
            else:
                x = (r.width - ancho_texto) / 2
            y = (margen if posicion.startswith('arriba') else r.height - margen)
            # rotate=pagina.rotation respeta páginas giradas
            pagina.insert_text((x, y), texto, fontsize=tamano, color=(0.25, 0.25, 0.25))
        buf = io.BytesIO()
        doc.save(buf, garbage=3, deflate=True)
        return buf.getvalue()
    finally:
        doc.close()


# ============================================================
# DIVIDIR POR RANGOS
# ============================================================

def dividir(contenido_pdf, rangos_texto):
    """Divide el PDF según rangos tipo "1-3, 5, 8-10".

    Returns:
        Lista de tuplas (nombre_sugerido, bytes_pdf), una por rango.
    """
    try:
        doc = fitz.open(stream=contenido_pdf, filetype='pdf')
    except Exception as e:
        raise ErrorConversion('PDF inválido: %s' % e)

    total = len(doc)
    partes = []
    try:
        for trozo in (rangos_texto or '').split(','):
            trozo = trozo.strip()
            if not trozo:
                continue
            if '-' in trozo:
                ini, fin = trozo.split('-', 1)
            else:
                ini = fin = trozo
            try:
                ini, fin = int(ini), int(fin)
            except ValueError:
                raise ErrorConversion('Rango inválido: "%s" (usa por ejemplo 1-3, 5, 8-10)' % trozo)
            if ini < 1 or fin > total or ini > fin:
                raise ErrorConversion('Rango "%s" fuera del documento (tiene %d páginas)' % (trozo, total))

            nuevo = fitz.open()
            nuevo.insert_pdf(doc, from_page=ini - 1, to_page=fin - 1)
            buf = io.BytesIO()
            nuevo.save(buf, garbage=3, deflate=True)
            nuevo.close()
            nombre = 'paginas_%d.pdf' % ini if ini == fin else 'paginas_%d-%d.pdf' % (ini, fin)
            partes.append((nombre, buf.getvalue()))
    finally:
        doc.close()

    if not partes:
        raise ErrorConversion('No se indicó ningún rango de páginas.')
    return partes


# ============================================================
# DESBLOQUEAR (quitar contraseña)
# ============================================================

def desbloquear(contenido_pdf, password):
    """Quita la contraseña de un PDF (requiere conocerla — no rompe claves)."""
    try:
        doc = fitz.open(stream=contenido_pdf, filetype='pdf')
    except Exception as e:
        raise ErrorConversion('PDF inválido: %s' % e)

    try:
        if doc.needs_pass:
            if not doc.authenticate(password or ''):
                raise ErrorConversion('Contraseña incorrecta.')
        elif not doc.is_encrypted:
            raise ErrorConversion('Este PDF no tiene contraseña.')
        buf = io.BytesIO()
        # encryption=PDF_ENCRYPT_NONE elimina toda la protección
        doc.save(buf, garbage=3, deflate=True, encryption=fitz.PDF_ENCRYPT_NONE)
        return buf.getvalue()
    finally:
        doc.close()
